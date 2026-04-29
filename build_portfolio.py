"""Generate Hanova Consultancy company portfolio (PPTX)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---------- Logo assets ----------
_BRAND_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "branding")
LOGO_FULL    = os.path.join(_BRAND_DIR, "hanova-logo-transparent.png")  # navy on light
LOGO_FULL_ON_DARK = os.path.join(_BRAND_DIR, "hanova-logo-onDark.png")  # white on dark
LOGO_MARK    = os.path.join(_BRAND_DIR, "hanova-mark.png")
LOGO_MARK_ON_DARK = os.path.join(_BRAND_DIR, "hanova-mark-onDark.png")

# ---------- Brand palette (from Hanova logo) ----------
INK        = RGBColor(0x14, 0x25, 0x4A)  # deep navy (logo "H" + wordmark)
INK_SOFT   = RGBColor(0x1B, 0x30, 0x5C)  # softer navy
GOLD       = RGBColor(0x14, 0xB3, 0x91)  # teal/emerald (logo arrow) — primary accent
GOLD_SOFT  = RGBColor(0x6F, 0xD8, 0xC1)  # light teal
ORANGE     = RGBColor(0xE9, 0x6A, 0x2A)  # logo "IMPACT" accent
PAPER      = RGBColor(0xF6, 0xF8, 0xFA)  # cool off-white
MUTED      = RGBColor(0x6B, 0x76, 0x88)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LINE       = RGBColor(0x2A, 0x3D, 0x5C)

# ---------- Presentation setup (16:9) ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
             font="Calibri", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_line(slide, x, y, w, h, color=GOLD, weight=1.5):
    ln = slide.shapes.add_connector(1, x, y, x + w, y + h)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def page_frame(slide, bg=PAPER, ink=INK, page_num=None, total=None,
               eyebrow="HANOVA CONSULTANCY", show_logo=True):
    """Common page chrome: background, header rule, logo, footer."""
    add_rect(slide, 0, 0, SW, SH, bg)
    # Top accent bar
    add_rect(slide, 0, 0, SW, Inches(0.05), GOLD)
    # Header eyebrow
    add_text(slide, Inches(0.6), Inches(0.25), Inches(8), Inches(0.3),
             eyebrow, size=9, bold=True, color=ink, spacing=1.0)
    # Logo mark (top right) — actual brand mark
    if show_logo:
        mark_path = LOGO_MARK_ON_DARK if bg == INK else LOGO_MARK
        mark_h = Inches(0.5)
        mark_w = Inches(0.5 * 464 / 539)  # preserve aspect ratio
        slide.shapes.add_picture(mark_path,
                                 SW - Inches(0.6) - mark_w,
                                 Inches(0.18),
                                 width=mark_w, height=mark_h)
    # Footer rule
    add_rect(slide, Inches(0.6), SH - Inches(0.5), SW - Inches(1.2),
             Emu(6350), RGBColor(0xCC, 0xD3, 0xDE))
    # Footer text
    add_text(slide, Inches(0.6), SH - Inches(0.4), Inches(8), Inches(0.3),
             "Hanova Consultancy  ·  From Strategy to Measurable Impact",
             size=8, color=MUTED)
    if page_num is not None:
        add_text(slide, SW - Inches(2.0), SH - Inches(0.4), Inches(1.4),
                 Inches(0.3),
                 f"{page_num:02d} / {total:02d}", size=8, color=MUTED,
                 align=PP_ALIGN.RIGHT)


def section_title(slide, eyebrow, title, color=INK, eyebrow_color=GOLD):
    add_text(slide, Inches(0.6), Inches(0.95), Inches(8), Inches(0.35),
             eyebrow, size=10, bold=True, color=eyebrow_color, spacing=1.0)
    add_text(slide, Inches(0.6), Inches(1.3), Inches(11.5), Inches(1.1),
             title, size=34, bold=False, color=color, font="Georgia",
             spacing=1.05)
    add_line(slide, Inches(0.6), Inches(2.35), Inches(0.6), 0,
             color=GOLD, weight=2.25)


# ---------- Slide builders ----------
TOTAL = 14  # will adjust if changed

def slide_cover():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    # Subtle right band
    add_rect(s, SW - Inches(4.2), 0, Inches(4.2), SH, INK_SOFT)
    # Teal accent
    add_rect(s, 0, 0, Inches(0.18), SH, GOLD)
    # Brand mark (on-dark variant) + wordmark
    mh = Inches(1.1)
    mw = Inches(1.1 * 464 / 539)
    s.shapes.add_picture(LOGO_MARK_ON_DARK, Inches(0.85), Inches(0.7),
                         width=mw, height=mh)
    add_text(s, Inches(0.85) + mw + Inches(0.25), Inches(0.85),
             Inches(6), Inches(0.5),
             "HANOVA", size=22, bold=True, color=WHITE, spacing=1.0)
    add_text(s, Inches(0.85) + mw + Inches(0.25), Inches(1.25),
             Inches(6), Inches(0.35),
             "C O N S U L T A N C Y", size=10, color=GOLD_SOFT, spacing=1.0)

    # Eyebrow
    add_text(s, Inches(0.85), Inches(2.6), Inches(10), Inches(0.4),
             "Strategy  ·  Transformation  ·  Results",
             size=11, bold=True, color=GOLD, spacing=1.0)

    # Headline
    add_text(s, Inches(0.85), Inches(3.0), Inches(11), Inches(1.4),
             "From Strategy to", size=54, color=WHITE,
             font="Georgia", spacing=1.0)
    add_text(s, Inches(0.85), Inches(3.85), Inches(11), Inches(1.4),
             "Measurable Impact.", size=54, color=GOLD,
             font="Georgia", italic=True, spacing=1.0)

    # Sub
    add_text(s, Inches(0.85), Inches(5.2), Inches(8.5), Inches(1.2),
             "Company Portfolio  ·  2026", size=12, color=WHITE, spacing=1.2)
    add_text(s, Inches(0.85), Inches(5.55), Inches(8.5), Inches(1.2),
             "An execution-driven consultancy partnering with executives\n"
             "to transform organizational strategy into quantifiable outcomes.",
             size=12, color=RGBColor(0xCF, 0xDB, 0xE8), spacing=1.4)

    # Bottom rule + meta
    add_rect(s, Inches(0.85), SH - Inches(1.0), Inches(2.0),
             Emu(9525), GOLD)
    add_text(s, Inches(0.85), SH - Inches(0.8), Inches(6), Inches(0.3),
             "RIYADH  ·  KINGDOM OF SAUDI ARABIA", size=9, bold=True,
             color=GOLD_SOFT, spacing=1.0)
    add_text(s, Inches(0.85), SH - Inches(0.55), Inches(6), Inches(0.3),
             "hanovaconsultancy.com", size=9, color=WHITE, spacing=1.0)


def slide_contents():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=2, total=TOTAL)
    section_title(s, "01", "Contents")

    items = [
        ("01", "About Hanova"),
        ("02", "Our Approach & Philosophy"),
        ("03", "Four Pillars"),
        ("04", "Service Practice Areas"),
        ("05", "Industries We Serve"),
        ("06", "Why Hanova"),
        ("07", "Regulatory & Operational Fluency"),
        ("08", "Measured Impact"),
        ("09", "Executive Workshops"),
        ("10", "Engagement Model"),
        ("11", "Leadership"),
        ("12", "Contact"),
    ]
    col_w = Inches(5.6)
    row_h = Inches(0.42)
    start_y = Inches(2.7)
    for i, (num, label) in enumerate(items):
        col = i // 6
        row = i % 6
        x = Inches(0.6) + col * (col_w + Inches(0.4))
        y = start_y + row * row_h
        add_text(s, x, y, Inches(0.6), row_h, num,
                 size=11, bold=True, color=GOLD, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.7), y, col_w - Inches(0.7), row_h, label,
                 size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE,
                 font="Georgia")
        add_rect(s, x, y + row_h - Emu(3175),
                 col_w - Inches(0.2), Emu(3175),
                 RGBColor(0xD9, 0xDF, 0xE8))


def slide_about():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=3, total=TOTAL)
    section_title(s, "02  ·  ABOUT", "An execution-driven consultancy\nbuilt for complex transformation.")

    # Body left
    add_text(s, Inches(0.6), Inches(2.85), Inches(7.4), Inches(0.6),
             "Hanova Consultancy brings physician-led strategic thinking and "
             "deep operational expertise to the most challenging transformation "
             "mandates facing today's executives.",
             size=13, color=INK, spacing=1.45)
    add_text(s, Inches(0.6), Inches(4.0), Inches(7.4), Inches(0.6),
             "We don't deliver reports — we deliver results. Our multi-entity "
             "experience spans healthcare systems, financial institutions, "
             "government bodies, and growth-stage enterprises — each engagement "
             "anchored in KPI governance and measurable ROI.",
             size=13, color=INK, spacing=1.45)

    # Stat stack right
    stats = [("15+", "Years of Executive Advisory"),
             ("40+", "Transformation Mandates"),
             ("8",   "Industries Served"),
             ("100%", "Execution-Driven Engagements")]
    sx = Inches(8.6)
    sy = Inches(2.85)
    cw = Inches(2.05)
    ch = Inches(1.85)
    gap = Inches(0.15)
    for i, (num, label) in enumerate(stats):
        cx = sx + (i % 2) * (cw + gap)
        cy = sy + (i // 2) * (ch + gap)
        bg = INK if i == 3 else WHITE
        fg_num = GOLD if i == 3 else INK
        fg_lab = WHITE if i == 3 else MUTED
        add_rect(s, cx, cy, cw, ch, bg, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_text(s, cx + Inches(0.2), cy + Inches(0.25), cw - Inches(0.4),
                 Inches(0.9), num, size=34, bold=True, color=fg_num,
                 font="Georgia", spacing=1.0)
        add_text(s, cx + Inches(0.2), cy + Inches(1.15), cw - Inches(0.4),
                 Inches(0.7), label, size=9, bold=True, color=fg_lab,
                 spacing=1.2)


def slide_approach():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=4, total=TOTAL)
    section_title(s, "03  ·  APPROACH", "From insight to embedded outcome.")

    steps = [
        ("DISCOVER", "Diagnostic immersion across leadership, operations, "
                     "and data — establishing the real baseline."),
        ("DESIGN",   "Evidence-based strategy, prioritized roadmap, and "
                     "KPI architecture co-built with executive sponsors."),
        ("DELIVER",  "Hands-on execution alongside your teams — "
                     "systems, workflows, governance, and capability."),
        ("EMBED",    "Performance governance, training, and handover — "
                     "outcomes sustained long after we leave."),
    ]
    x0 = Inches(0.6)
    y0 = Inches(2.95)
    cw = Inches(2.95)
    ch = Inches(3.4)
    gap = Inches(0.18)
    for i, (k, v) in enumerate(steps):
        cx = x0 + i * (cw + gap)
        # card
        add_rect(s, cx, y0, cw, ch, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        # number
        add_text(s, cx + Inches(0.3), y0 + Inches(0.3), cw, Inches(0.5),
                 f"0{i+1}", size=11, bold=True, color=GOLD)
        # title
        add_text(s, cx + Inches(0.3), y0 + Inches(0.75), cw - Inches(0.6),
                 Inches(0.6), k, size=18, bold=True, color=INK,
                 font="Georgia", spacing=1.0)
        # rule
        add_rect(s, cx + Inches(0.3), y0 + Inches(1.4),
                 Inches(0.5), Emu(9525), GOLD)
        # body
        add_text(s, cx + Inches(0.3), y0 + Inches(1.65), cw - Inches(0.6),
                 ch - Inches(2.0), v, size=11, color=INK, spacing=1.45)


def slide_pillars():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, SW, Inches(0.05), GOLD)
    add_text(s, Inches(0.6), Inches(0.25), Inches(8), Inches(0.3),
             "HANOVA CONSULTANCY", size=9, bold=True, color=WHITE)
    _mh = Inches(0.5); _mw = Inches(0.5 * 464 / 539)
    s.shapes.add_picture(LOGO_MARK_ON_DARK, SW - Inches(0.6) - _mw, Inches(0.18), width=_mw, height=_mh)
    add_text(s, Inches(0.6), Inches(0.95), Inches(8), Inches(0.35),
             "04  ·  FOUNDATION", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(1.3), Inches(11.5), Inches(1.1),
             "Four pillars. One outcome.", size=34, color=WHITE,
             font="Georgia", spacing=1.05)
    add_rect(s, Inches(0.6), Inches(2.35), Inches(0.6), Emu(28575), GOLD)

    pillars = [
        ("◈", "Strategy",
         "Evidence-based strategic frameworks designed for real-world "
         "complexity and executive decision-making."),
        ("◉", "Data & Intelligence",
         "From raw data to boardroom-ready insight — BI architecture, "
         "KPI governance, and analytics maturity."),
        ("◎", "Digital Transformation",
         "End-to-end digital enablement — systems, integration, and "
         "scalable solutions that last beyond the engagement."),
        ("◆", "Measurable Results",
         "Every mandate is anchored in clear KPIs. We govern "
         "performance, not just process."),
    ]
    x0 = Inches(0.6)
    y0 = Inches(2.95)
    cw = Inches(2.95)
    ch = Inches(3.6)
    gap = Inches(0.18)
    for i, (ic, t, d) in enumerate(pillars):
        cx = x0 + i * (cw + gap)
        add_rect(s, cx, y0, cw, ch, INK_SOFT, line=LINE)
        add_text(s, cx + Inches(0.3), y0 + Inches(0.35), cw, Inches(0.7),
                 ic, size=30, color=GOLD, font="Georgia")
        add_text(s, cx + Inches(0.3), y0 + Inches(1.2), cw - Inches(0.6),
                 Inches(0.6), t, size=18, bold=True, color=WHITE,
                 font="Georgia")
        add_rect(s, cx + Inches(0.3), y0 + Inches(1.85),
                 Inches(0.5), Emu(9525), GOLD)
        add_text(s, cx + Inches(0.3), y0 + Inches(2.1), cw - Inches(0.6),
                 ch - Inches(2.4), d, size=11,
                 color=RGBColor(0xCF, 0xDB, 0xE8), spacing=1.45)

    add_rect(s, Inches(0.6), SH - Inches(0.5), SW - Inches(1.2),
             Emu(6350), RGBColor(0x2A, 0x3D, 0x5C))
    add_text(s, Inches(0.6), SH - Inches(0.4), Inches(8), Inches(0.3),
             "Hanova Consultancy  ·  From Strategy to Measurable Impact",
             size=8, color=GOLD_SOFT)
    add_text(s, SW - Inches(2.0), SH - Inches(0.4), Inches(1.4), Inches(0.3),
             f"05 / {TOTAL:02d}", size=8, color=GOLD_SOFT, align=PP_ALIGN.RIGHT)


def slide_services():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=6, total=TOTAL)
    section_title(s, "05  ·  PRACTICE AREAS", "What we do.")

    services = [
        ("01", "Strategy & Transformation",
         "Corporate strategy, organizational redesign, and full-cycle "
         "transformation programs."),
        ("02", "Operations & Revenue Optimization",
         "Workflow redesign, revenue cycle management, and operational "
         "efficiency for sustainable growth."),
        ("03", "Data & Business Intelligence",
         "End-to-end BI implementation, data governance, and real-time "
         "analytics infrastructure."),
        ("04", "Digital Transformation",
         "Technology roadmaps, system modernization, and digital "
         "capability building at enterprise scale."),
        ("05", "Integration & Interoperability",
         "System integration architecture, API ecosystems, and "
         "cross-platform data flows."),
        ("06", "Custom Digital Solutions",
         "Bespoke software, intelligent automation, and purpose-built "
         "platforms for complex needs."),
    ]
    cols = 3
    rows = 2
    x0 = Inches(0.6)
    y0 = Inches(2.85)
    cw = Inches(4.05)
    ch = Inches(2.05)
    gx = Inches(0.15)
    gy = Inches(0.18)
    for i, (n, t, d) in enumerate(services):
        r, c = divmod(i, cols)
        cx = x0 + c * (cw + gx)
        cy = y0 + r * (ch + gy)
        add_rect(s, cx, cy, cw, ch, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_rect(s, cx, cy, Inches(0.05), ch, GOLD)
        add_text(s, cx + Inches(0.3), cy + Inches(0.2), Inches(1), Inches(0.4),
                 n, size=11, bold=True, color=GOLD)
        add_text(s, cx + Inches(0.3), cy + Inches(0.55), cw - Inches(0.5),
                 Inches(0.65), t, size=14, bold=True, color=INK,
                 font="Georgia", spacing=1.1)
        add_text(s, cx + Inches(0.3), cy + Inches(1.2), cw - Inches(0.5),
                 ch - Inches(1.3), d, size=10.5, color=INK_SOFT, spacing=1.4)


def slide_industries():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=7, total=TOTAL)
    section_title(s, "06  ·  INDUSTRIES", "Where we work.")

    inds = [
        ("Healthcare & Life Sciences",
         "Hospitals, MOH entities, payers, and clinical networks."),
        ("Financial Services & Insurance",
         "Banks, insurers, and capital markets transformation."),
        ("Government & Public Sector",
         "Vision 2030 program delivery and public service modernization."),
        ("Retail & E-Commerce",
         "Omnichannel operations, customer data, and unit economics."),
        ("Logistics & Supply Chain",
         "Network design, visibility, and operational excellence."),
        ("Technology & Startups",
         "Scale-up advisory, GTM, and operating model design."),
    ]
    x0 = Inches(0.6)
    y0 = Inches(2.85)
    rw = SW - Inches(1.2)
    rh = Inches(0.62)
    gap = Inches(0.08)
    for i, (n, d) in enumerate(inds):
        ry = y0 + i * (rh + gap)
        add_rect(s, x0, ry, rw, rh, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_text(s, x0 + Inches(0.3), ry, Inches(0.6), rh,
                 f"0{i+1}", size=11, bold=True, color=GOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x0 + Inches(1.0), ry, Inches(4.6), rh,
                 n, size=14, bold=True, color=INK, font="Georgia",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x0 + Inches(5.7), ry, rw - Inches(6.0), rh,
                 d, size=11, color=INK_SOFT, anchor=MSO_ANCHOR.MIDDLE)


def slide_why():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=8, total=TOTAL)
    section_title(s, "07  ·  WHY HANOVA",
                  "Differentiated by design.\nProven in practice.")

    items = [
        ("Physician-Led Transformation",
         "Clinical leadership background enables precision in highly "
         "regulated, complex systems."),
        ("Real Multi-Entity Experience",
         "We have operated inside the systems we advise — across "
         "hospitals, government, finance, and tech."),
        ("Execution-Driven, Not Report-Driven",
         "Our mandate ends when results are embedded — not when the "
         "deck is delivered."),
        ("KPI Governance & Revenue Expertise",
         "We build performance governance that outlasts the engagement."),
        ("Regulatory & Operational Fluency",
         "Deep familiarity with NDMO, PDPL, JCI, CBAHI, and Vision 2030."),
        ("Senior-Led, Every Engagement",
         "No junior associates managing your transformation. "
         "Senior principals own delivery."),
    ]
    cols = 2
    x0 = Inches(0.6)
    y0 = Inches(2.85)
    cw = Inches(6.05)
    ch = Inches(1.4)
    gx = Inches(0.15)
    gy = Inches(0.15)
    for i, (t, d) in enumerate(items):
        r, c = divmod(i, cols)
        cx = x0 + c * (cw + gx)
        cy = y0 + r * (ch + gy)
        add_rect(s, cx, cy, cw, ch, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_rect(s, cx + Inches(0.3), cy + Inches(0.35), Inches(0.25),
                 Emu(19050), GOLD)
        add_text(s, cx + Inches(0.7), cy + Inches(0.2), cw - Inches(1),
                 Inches(0.45), t, size=13, bold=True, color=INK,
                 font="Georgia")
        add_text(s, cx + Inches(0.7), cy + Inches(0.65), cw - Inches(1),
                 ch - Inches(0.8), d, size=10.5, color=INK_SOFT, spacing=1.4)


def slide_regulatory():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=9, total=TOTAL)
    section_title(s, "08  ·  COMPLIANCE",
                  "Regulatory & operational fluency.")

    add_text(s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.6),
             "We operate fluently within the frameworks that govern our "
             "clients' industries — accelerating delivery and de-risking "
             "transformation.",
             size=13, color=INK_SOFT, spacing=1.4)

    badges = [
        ("NDMO", "National Data Management Office"),
        ("PDPL", "Personal Data Protection Law"),
        ("JCI", "Joint Commission International"),
        ("CBAHI", "Saudi Central Board for Accreditation"),
        ("Vision 2030", "Saudi Arabia transformation programs"),
        ("HIPAA-aligned", "Health information governance"),
    ]
    x0 = Inches(0.6)
    y0 = Inches(4.1)
    cw = Inches(2.02)
    ch = Inches(2.0)
    gap = Inches(0.13)
    for i, (k, sub) in enumerate(badges):
        cx = x0 + i * (cw + gap)
        add_rect(s, cx, y0, cw, ch, INK)
        add_text(s, cx, y0 + Inches(0.45), cw, Inches(0.7),
                 k, size=18, bold=True, color=GOLD,
                 align=PP_ALIGN.CENTER, font="Georgia")
        add_rect(s, cx + cw/2 - Inches(0.25), y0 + Inches(1.15),
                 Inches(0.5), Emu(9525), GOLD)
        add_text(s, cx + Inches(0.15), y0 + Inches(1.3),
                 cw - Inches(0.3), Inches(0.65),
                 sub, size=9, color=WHITE,
                 align=PP_ALIGN.CENTER, spacing=1.3)


def slide_impact():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, SW, Inches(0.05), GOLD)
    add_text(s, Inches(0.6), Inches(0.25), Inches(8), Inches(0.3),
             "HANOVA CONSULTANCY", size=9, bold=True, color=WHITE)
    _mh = Inches(0.5); _mw = Inches(0.5 * 464 / 539)
    s.shapes.add_picture(LOGO_MARK_ON_DARK, SW - Inches(0.6) - _mw, Inches(0.18), width=_mw, height=_mh)
    add_text(s, Inches(0.6), Inches(0.95), Inches(8), Inches(0.35),
             "09  ·  MEASURED IMPACT", size=10, bold=True, color=GOLD)
    add_text(s, Inches(0.6), Inches(1.3), Inches(11.5), Inches(1.1),
             "Results our clients have achieved.",
             size=34, color=WHITE, font="Georgia", spacing=1.05)
    add_rect(s, Inches(0.6), Inches(2.35), Inches(0.6), Emu(28575), GOLD)

    impact = [
        ("+32%", "Revenue Cycle Efficiency",
         "Healthcare system, post-operational redesign"),
        ("-45%", "Operational Overhead",
         "Government entity, process optimization"),
        ("3×",  "Data Decision Speed",
         "Financial group, BI transformation"),
        ("+28%", "Patient Satisfaction",
         "Multi-specialty hospital, journey redesign"),
    ]
    x0 = Inches(0.6)
    y0 = Inches(2.95)
    cw = Inches(2.95)
    ch = Inches(3.4)
    gap = Inches(0.18)
    for i, (m, a, d) in enumerate(impact):
        cx = x0 + i * (cw + gap)
        add_rect(s, cx, y0, cw, ch, INK_SOFT, line=LINE)
        add_text(s, cx + Inches(0.3), y0 + Inches(0.4), cw - Inches(0.6),
                 Inches(1.4), m, size=48, bold=True, color=GOLD,
                 font="Georgia", spacing=1.0)
        add_rect(s, cx + Inches(0.3), y0 + Inches(1.85),
                 Inches(0.5), Emu(9525), GOLD)
        add_text(s, cx + Inches(0.3), y0 + Inches(2.05), cw - Inches(0.6),
                 Inches(0.6), a, size=14, bold=True, color=WHITE,
                 font="Georgia")
        add_text(s, cx + Inches(0.3), y0 + Inches(2.65), cw - Inches(0.6),
                 ch - Inches(2.8), d, size=10,
                 color=RGBColor(0xCF, 0xDB, 0xE8), spacing=1.4)

    add_rect(s, Inches(0.6), SH - Inches(0.5), SW - Inches(1.2),
             Emu(6350), RGBColor(0x2A, 0x3D, 0x5C))
    add_text(s, Inches(0.6), SH - Inches(0.4), Inches(8), Inches(0.3),
             "Selected outcomes — full case studies available on request.",
             size=8, color=GOLD_SOFT)
    add_text(s, SW - Inches(2.0), SH - Inches(0.4), Inches(1.4), Inches(0.3),
             f"10 / {TOTAL:02d}", size=8, color=GOLD_SOFT, align=PP_ALIGN.RIGHT)


def slide_workshops():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=11, total=TOTAL)
    section_title(s, "10  ·  EXECUTIVE EDUCATION",
                  "Workshops designed for leaders who act.")

    add_text(s, Inches(0.6), Inches(2.85), Inches(6.4), Inches(2.5),
             "Our executive workshops are intensive, outcome-oriented "
             "programs that arm leadership teams with the frameworks, "
             "tools, and clarity required for transformational "
             "decision-making.\n\n"
             "Delivered on-site or in immersive off-site formats, each "
             "engagement is custom-scoped to the client's strategic "
             "priorities and KPI mandates.",
             size=12, color=INK, spacing=1.5)

    sessions = [
        "Healthcare Transformation for Executives",
        "Data-Driven Decision Making",
        "KPI Governance & Performance Management",
        "Revenue & Operations Optimization",
        "Digital Transformation for Leaders",
    ]
    x = Inches(7.5)
    y = Inches(2.85)
    w = Inches(5.2)
    h = Inches(0.65)
    gap = Inches(0.1)
    for i, t in enumerate(sessions):
        ry = y + i * (h + gap)
        add_rect(s, x, ry, w, h, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_rect(s, x, ry, Inches(0.05), h, GOLD)
        add_text(s, x + Inches(0.3), ry, Inches(0.5), h,
                 f"0{i+1}", size=10, bold=True, color=GOLD,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x + Inches(0.85), ry, w - Inches(1), h,
                 t, size=12, bold=True, color=INK, font="Georgia",
                 anchor=MSO_ANCHOR.MIDDLE)


def slide_engagement():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=12, total=TOTAL)
    section_title(s, "11  ·  ENGAGEMENT MODEL",
                  "How we work with you.")

    rows = [
        ("Advisory Retainer",
         "Ongoing senior advisory to the executive office.",
         "Quarterly  ·  Renewable"),
        ("Transformation Mandate",
         "End-to-end program delivery with embedded teams.",
         "6 – 18 months"),
        ("Diagnostic Sprint",
         "Fixed-scope diagnostic with prioritized roadmap.",
         "4 – 8 weeks"),
        ("Executive Workshop",
         "Custom-scoped leadership program, on or off-site.",
         "1 – 5 days"),
    ]
    x0 = Inches(0.6)
    y0 = Inches(2.85)
    rw = SW - Inches(1.2)
    rh = Inches(0.85)
    gap = Inches(0.12)

    # header
    hdr_y = y0
    add_rect(s, x0, hdr_y, rw, Inches(0.45), INK)
    add_text(s, x0 + Inches(0.3), hdr_y, Inches(3.5), Inches(0.45),
             "FORMAT", size=9, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0 + Inches(4.0), hdr_y, Inches(5.5), Inches(0.45),
             "DESCRIPTION", size=9, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x0 + Inches(9.7), hdr_y, Inches(2.4), Inches(0.45),
             "TYPICAL DURATION", size=9, bold=True, color=GOLD,
             anchor=MSO_ANCHOR.MIDDLE)

    y = hdr_y + Inches(0.45) + gap
    for i, (a, b, c) in enumerate(rows):
        add_rect(s, x0, y, rw, rh, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_rect(s, x0, y, Inches(0.05), rh, GOLD)
        add_text(s, x0 + Inches(0.3), y, Inches(3.6), rh,
                 a, size=14, bold=True, color=INK, font="Georgia",
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x0 + Inches(4.0), y, Inches(5.5), rh,
                 b, size=11, color=INK_SOFT,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, x0 + Inches(9.7), y, Inches(2.4), rh,
                 c, size=11, bold=True, color=INK,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += rh + gap


def slide_leadership():
    s = prs.slides.add_slide(BLANK)
    page_frame(s, page_num=13, total=TOTAL)
    section_title(s, "12  ·  LEADERSHIP",
                  "Senior-led. Always.")

    add_text(s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.6),
             "Every Hanova engagement is owned and delivered by senior "
             "principals — not staffed through junior associates. Our "
             "leadership combines clinical, operational, and technology "
             "depth with C-suite advisory experience.",
             size=12, color=INK_SOFT, spacing=1.5)

    creds = [
        ("Physician Leadership", "Clinical insight, regulated systems"),
        ("Executive Advisory",   "C-suite and board-level partnership"),
        ("Operational Depth",    "Built and run the systems we advise"),
        ("Technology Fluency",   "BI, data, integration, automation"),
        ("Public Sector Experience", "Vision 2030 program delivery"),
        ("Regulatory Mastery",   "NDMO  ·  PDPL  ·  JCI  ·  CBAHI"),
    ]
    cols = 3
    x0 = Inches(0.6)
    y0 = Inches(4.0)
    cw = Inches(4.05)
    ch = Inches(1.25)
    gx = Inches(0.13)
    gy = Inches(0.13)
    for i, (t, d) in enumerate(creds):
        r, c = divmod(i, cols)
        cx = x0 + c * (cw + gx)
        cy = y0 + r * (ch + gy)
        add_rect(s, cx, cy, cw, ch, WHITE, line=RGBColor(0xD9, 0xDF, 0xE8))
        add_rect(s, cx, cy, Inches(0.05), ch, GOLD)
        add_text(s, cx + Inches(0.3), cy + Inches(0.2), cw - Inches(0.5),
                 Inches(0.45), t, size=13, bold=True, color=INK,
                 font="Georgia")
        add_text(s, cx + Inches(0.3), cy + Inches(0.65), cw - Inches(0.5),
                 ch - Inches(0.75), d, size=10.5, color=INK_SOFT, spacing=1.3)


def slide_contact():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, 0, 0, SW, SH, INK)
    add_rect(s, 0, 0, Inches(0.18), SH, GOLD)
    add_rect(s, SW - Inches(4.5), 0, Inches(4.5), SH, INK_SOFT)
    # Brand mark on right panel (top)
    cmh = Inches(1.4)
    cmw = Inches(1.4 * 464 / 539)
    s.shapes.add_picture(LOGO_MARK_ON_DARK,
                         SW - Inches(4.5) + (Inches(4.5) - cmw) / 2,
                         Inches(0.65),
                         width=cmw, height=cmh)

    add_text(s, Inches(0.85), Inches(0.9), Inches(8), Inches(0.4),
             "13  ·  CONTACT", size=11, bold=True, color=GOLD)
    add_text(s, Inches(0.85), Inches(1.4), Inches(11), Inches(1.6),
             "Let's talk transformation.",
             size=46, color=WHITE, font="Georgia", spacing=1.05)
    add_text(s, Inches(0.85), Inches(2.6), Inches(11), Inches(1.6),
             "Speak directly with a senior consultant.",
             size=22, color=GOLD, font="Georgia", italic=True, spacing=1.1)
    add_rect(s, Inches(0.85), Inches(3.6), Inches(0.6), Emu(28575), GOLD)

    add_text(s, Inches(0.85), Inches(3.9), Inches(7), Inches(0.6),
             "No intake forms. No junior staff. Just a focused "
             "conversation about your organization's most pressing "
             "challenges.",
             size=13, color=RGBColor(0xCF, 0xDB, 0xE8), spacing=1.5)

    # Right column contact details
    cx = SW - Inches(4.2)
    cy = Inches(2.4)
    add_text(s, cx, cy, Inches(3.7), Inches(0.4),
             "DIRECT", size=10, bold=True, color=GOLD)
    add_text(s, cx, cy + Inches(0.4), Inches(3.7), Inches(0.5),
             "info@hanovaconsultancy.com", size=14, color=WHITE)

    add_text(s, cx, cy + Inches(1.2), Inches(3.7), Inches(0.4),
             "WEB", size=10, bold=True, color=GOLD)
    add_text(s, cx, cy + Inches(1.6), Inches(3.7), Inches(0.5),
             "hanovaconsultancy.com", size=14, color=WHITE)

    add_text(s, cx, cy + Inches(2.4), Inches(3.7), Inches(0.4),
             "OFFICE", size=10, bold=True, color=GOLD)
    add_text(s, cx, cy + Inches(2.8), Inches(3.7), Inches(0.6),
             "Riyadh\nKingdom of Saudi Arabia",
             size=14, color=WHITE, spacing=1.4)

    add_rect(s, Inches(0.85), SH - Inches(0.9), Inches(2.0),
             Emu(9525), GOLD)
    add_text(s, Inches(0.85), SH - Inches(0.7), Inches(8), Inches(0.3),
             "HANOVA  ·  CONSULTANCY",
             size=10, bold=True, color=WHITE)
    add_text(s, Inches(0.85), SH - Inches(0.45), Inches(8), Inches(0.3),
             "From Strategy to Measurable Impact.",
             size=9, color=GOLD_SOFT, italic=True)


# ---------- Build deck ----------
slide_cover()        # 1
slide_contents()     # 2
slide_about()        # 3
slide_approach()     # 4
slide_pillars()      # 5
slide_services()     # 6
slide_industries()   # 7
slide_why()          # 8
slide_regulatory()   # 9
slide_impact()       # 10
slide_workshops()    # 11
slide_engagement()   # 12
slide_leadership()   # 13
slide_contact()      # 14

import os
base = r"d:\Hanova\Hanova-Consultancy-Portfolio.pptx"
out = base
i = 1
while True:
    try:
        prs.save(out)
        break
    except PermissionError:
        out = base.replace(".pptx", f"-v{i}.pptx")
        i += 1
print(f"Saved: {out}  ·  {len(prs.slides)} slides")
